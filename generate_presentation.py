"""
Generate scTRaCT PowerPoint presentation.
Run: python generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ─────────────────────────────────────────────
# Color palette
# ─────────────────────────────────────────────
C_DARK_BLUE  = RGBColor(0x1A, 0x3A, 0x5C)   # headers / accent
C_MID_BLUE   = RGBColor(0x2E, 0x6D, 0xA4)   # sub-headers
C_LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF5)   # subtle bg boxes
C_TEAL       = RGBColor(0x1A, 0x8C, 0x8C)   # highlight
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY       = RGBColor(0x55, 0x55, 0x77)
C_GREEN      = RGBColor(0x27, 0xAE, 0x60)
C_ORANGE     = RGBColor(0xE6, 0x7E, 0x22)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # blank


# ─────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE if False else 1,  # MSO_SHAPE.RECTANGLE = 1
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, italic=False,
                 color=C_NEAR_BLACK, align=PP_ALIGN.LEFT,
                 word_wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_bullet_slide_content(slide, bullets, x, y, w, h,
                              font_size=18, color=C_NEAR_BLACK,
                              bullet_indent=0.3):
    """Add multiple bullet points with optional sub-bullets (list of str or (str, list))."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for item in bullets:
        if isinstance(item, tuple):
            main_text, sub_items = item
        else:
            main_text, sub_items = item, []

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.level = 0
        run = p.add_run()
        run.text = main_text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"

        for sub in sub_items:
            sp = tf.add_paragraph()
            sp.level = 1
            sr = sp.add_run()
            sr.text = sub
            sr.font.size = Pt(font_size - 2)
            sr.font.color.rgb = C_GRAY
            sr.font.name = "Calibri"

    return txBox


def add_section_header(slide, title, subtitle=None):
    """Top banner with dark-blue background."""
    add_rect(slide, 0, 0, 13.33, 1.3, fill_color=C_DARK_BLUE)
    add_text_box(slide, title, 0.3, 0.1, 12.5, 0.75,
                 font_size=32, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.3, 0.78, 12.5, 0.45,
                     font_size=16, italic=True, color=RGBColor(0xCC, 0xDD, 0xFF),
                     align=PP_ALIGN.LEFT)


def add_slide_number(slide, num):
    add_text_box(slide, str(num), 12.9, 7.2, 0.4, 0.3,
                 font_size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK_BLUE)

# Decorative top bar
add_rect(slide, 0, 0, 13.33, 0.25, fill_color=C_TEAL)

# Main title
add_text_box(slide, "scTRaCT", 0.5, 1.0, 12.33, 1.3,
             font_size=60, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide,
             "Single-Cell Transformer for Cell Type Classification\n"
             "with Explainable AI Gene Attribution",
             0.5, 2.3, 12.33, 1.2,
             font_size=22, color=RGBColor(0xAA, 0xCC, 0xFF),
             align=PP_ALIGN.CENTER)

# Light divider
add_rect(slide, 2.5, 3.55, 8.33, 0.04, fill_color=C_TEAL)

# Author / affiliation block
add_text_box(slide,
             "PhD Research Presentation\nUniversity of Texas at San Antonio",
             0.5, 3.7, 12.33, 0.9,
             font_size=16, italic=True, color=RGBColor(0xCC, 0xDD, 0xFF),
             align=PP_ALIGN.CENTER)

# Tech stack badge-like strip
add_rect(slide, 0, 6.8, 13.33, 0.7, fill_color=RGBColor(0x12, 0x28, 0x40))
add_text_box(slide,
             "PyTorch  •  Captum  •  Scanpy / AnnData  •  Integrated Gradients  •  GradientSHAP",
             0.5, 6.85, 12.33, 0.55,
             font_size=13, color=RGBColor(0x88, 0xBB, 0xFF),
             align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 2 — Outline
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "Presentation Outline")
add_slide_number(slide, 2)

items = [
    "1.  The Biological Problem — why cell type classification matters",
    "2.  scTRaCT Architecture Overview",
    "3.  Data Inputs: Log-Normalized Counts & MCA Distances",
    "4.  Transformer Model Deep Dive",
    "5.  Training: Focal Loss & Class Imbalance",
    "6.  What is Explainable AI (XAI)?",
    "7.  Integrated Gradients — Theory & Implementation",
    "8.  GradientSHAP — Theory & Implementation",
    "9.  scTRaCT XAI Pipeline (Step-by-Step)",
    "10. Outputs — Bar Charts, Heatmaps, CSVs",
    "11. Summary & Future Work",
]

add_bullet_slide_content(slide, items, 1.0, 1.5, 11.5, 5.8,
                         font_size=17, color=C_NEAR_BLACK)


# ─────────────────────────────────────────────
# SLIDE 3 — Biological Motivation
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "The Biological Problem",
                   "Why do we need to classify single cells?")
add_slide_number(slide, 3)

left_bullets = [
    ("Single-cell RNA-seq (scRNA-seq) measures gene expression in thousands of individual cells simultaneously.",
     []),
    ("Each cell has a unique gene expression profile — its molecular fingerprint.", []),
    ("Cell type identity drives tissue function, disease progression, and therapy response.", []),
    ("Manual annotation by domain experts is slow, subjective, and does not scale to millions of cells.", []),
    ("Automated, accurate, and interpretable classifiers are therefore critical for translational research.", []),
]

add_bullet_slide_content(slide, left_bullets, 0.5, 1.45, 12.5, 5.5,
                         font_size=17, color=C_NEAR_BLACK)

# Side highlight box
add_rect(slide, 9.3, 1.55, 3.7, 2.4, fill_color=C_LIGHT_BLUE,
         line_color=C_MID_BLUE, line_width_pt=1.5)
add_text_box(slide, "Key Challenge", 9.45, 1.65, 3.4, 0.45,
             font_size=13, bold=True, color=C_DARK_BLUE)
add_text_box(slide,
             "Rare cell types are underrepresented in training data, causing standard classifiers to ignore them.",
             9.45, 2.05, 3.4, 1.8, font_size=13, color=C_NEAR_BLACK, word_wrap=True)


# ─────────────────────────────────────────────
# SLIDE 4 — Architecture Overview
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "scTRaCT Architecture Overview")
add_slide_number(slide, 4)

# Two-column layout
# LEFT: flow boxes
boxes = [
    (0.4, 1.6,  3.2, 0.6, C_LIGHT_BLUE, C_MID_BLUE, "Input: scRNA-seq AnnData (.h5ad)"),
    (0.4, 2.45, 3.2, 0.6, C_LIGHT_BLUE, C_MID_BLUE, "Preprocessing: log-norm + MCA distances"),
    (0.4, 3.3,  3.2, 0.6, C_LIGHT_BLUE, C_MID_BLUE, "Two Linear Embeddings (count + distance)"),
    (0.4, 4.15, 3.2, 0.6, C_LIGHT_BLUE, C_MID_BLUE, "CLS Token Prepended → 3-Layer Transformer"),
    (0.4, 5.0,  3.2, 0.6, C_LIGHT_BLUE, C_MID_BLUE, "CLS Output → Fully Connected → Softmax"),
    (0.4, 5.85, 3.2, 0.6, RGBColor(0xD5, 0xF5, 0xE3), C_GREEN, "Predicted Cell Type  ✓"),
]
for bx, by, bw, bh, fc, lc, bt in boxes:
    add_rect(slide, bx, by, bw, bh, fill_color=fc, line_color=lc, line_width_pt=1.5)
    add_text_box(slide, bt, bx + 0.1, by + 0.1, bw - 0.2, bh - 0.1,
                 font_size=13, color=C_NEAR_BLACK, bold=False)

# Arrows between boxes (simple text arrows)
for ay in [2.25, 3.1, 3.95, 4.8, 5.65]:
    add_text_box(slide, "▼", 1.85, ay, 0.4, 0.25, font_size=14, color=C_MID_BLUE, align=PP_ALIGN.CENTER)

# RIGHT: key facts
add_rect(slide, 4.2, 1.5, 8.8, 5.6, fill_color=C_LIGHT_BLUE,
         line_color=C_MID_BLUE, line_width_pt=1)
add_text_box(slide, "Key Design Decisions", 4.4, 1.6, 8.4, 0.45,
             font_size=15, bold=True, color=C_DARK_BLUE)

facts = [
    ("Dual Input Streams", [
        "Count embeddings encode gene expression magnitude",
        "Distance embeddings encode MCA cell-to-gene distances",
        "Both are projected to the same embedding dimension, then concatenated",
    ]),
    ("CLS Token (borrowed from BERT)", [
        "A learnable [CLS] vector is prepended to the sequence",
        "After all attention layers, its representation summarizes the entire cell",
        "Only the CLS position is passed to the classifier head",
    ]),
    ("3-Layer Multi-Head Self-Attention (8 heads)", [
        "Each layer: Attention → Residual + LayerNorm",
        "Final layer adds a Feed-Forward Network (FFN) block",
    ]),
    ("FocalLoss — handles class imbalance", [
        "Down-weights easy examples so rare cell types get more gradient",
    ]),
]

add_bullet_slide_content(slide, facts, 4.35, 2.1, 8.5, 4.8,
                         font_size=13, color=C_NEAR_BLACK)


# ─────────────────────────────────────────────
# SLIDE 5 — Data Inputs
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "Data Inputs: Two Complementary Views of Each Cell")
add_slide_number(slide, 5)

# Stream 1
add_rect(slide, 0.3, 1.5, 5.9, 5.6, fill_color=C_LIGHT_BLUE,
         line_color=C_MID_BLUE, line_width_pt=1.5)
add_text_box(slide, "Stream 1 — Log-Normalized Counts",
             0.5, 1.6, 5.5, 0.5, font_size=15, bold=True, color=C_DARK_BLUE)
add_bullet_slide_content(slide, [
    "Raw UMI counts are log1p-normalized per cell",
    "Layer key: adata.layers['lognorm']",
    "Shape: (cells × genes) sparse matrix",
    "Captures which genes are expressed and by how much",
    "Standard preprocessing in Scanpy / scRNA-seq pipelines",
    "Baseline for IG = zero vector (silent cell)",
], 0.5, 2.15, 5.5, 4.7, font_size=14, color=C_NEAR_BLACK)

# Stream 2
add_rect(slide, 6.8, 1.5, 6.2, 5.6, fill_color=RGBColor(0xFD, 0xF2, 0xE9),
         line_color=C_ORANGE, line_width_pt=1.5)
add_text_box(slide, "Stream 2 — MCA Distance Features",
             7.0, 1.6, 5.8, 0.5, font_size=15, bold=True, color=C_ORANGE)
add_bullet_slide_content(slide, [
    "Multiple Correspondence Analysis (MCA) maps cells and genes into a shared factor space",
    "Euclidean distances in this space = gene-to-cell affinity",
    "Layer key: adata.layers['distance_matrix']",
    "Transformed to inverse distance:  1 / (d + ε)  before embedding",
    "This gives high weight to nearby (highly associated) genes",
    "Adds structural, geometry-aware information beyond raw counts",
], 7.0, 2.15, 5.8, 4.7, font_size=14, color=C_NEAR_BLACK)

add_text_box(slide, "⟷  Both streams are embedded independently, then concatenated before the Transformer  ⟷",
             0.3, 7.05, 12.73, 0.35, font_size=12, italic=True,
             color=C_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 6 — Transformer Deep Dive
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "Transformer Model — Forward Pass",
                   "model.py  |  TransformerModel")
add_slide_number(slide, 6)

steps = [
    ("Step 1 — Embed", C_LIGHT_BLUE, C_MID_BLUE,
     "count_embed = Linear(num_genes → 1024)(x_counts)\n"
     "dist_embed  = Linear(num_genes → 1024)(x_dist)\n"
     "x = concat([count_embed, dist_embed], dim=-1)   # → dim 2048\n"
     "x = x.unsqueeze(1)   # shape: (batch, 1, 2048)"),

    ("Step 2 — CLS Token", C_LIGHT_BLUE, C_MID_BLUE,
     "cls = learnable parameter (1, 1, 2048)\n"
     "x = concat([cls, x], dim=1)   # shape: (batch, 2, 2048)\n"
     "x = LayerNorm(x)"),

    ("Step 3 — 3× Self-Attention Block", C_LIGHT_BLUE, C_MID_BLUE,
     "for each layer:\n"
     "  attn_out = MultiHeadAttention(x)   # 8 heads\n"
     "  x = x + Dropout(LayerNorm(attn_out))   # residual"),

    ("Step 4 — Feed-Forward Network", C_LIGHT_BLUE, C_MID_BLUE,
     "ffn_out = Linear(2048→2048)(GELU(Linear(2048→2048)(x)))\n"
     "x = x + Dropout(LayerNorm(ffn_out))   # residual"),

    ("Step 5 — Classify", RGBColor(0xD5, 0xF5, 0xE3), C_GREEN,
     "logits = Linear(2048 → num_classes)(x[:, 0])   # CLS position only\n"
     "prediction = argmax(logits)"),
]

y_pos = 1.5
for title, fc, lc, code in steps:
    add_rect(slide, 0.3, y_pos, 12.73, 0.95, fill_color=fc,
             line_color=lc, line_width_pt=1)
    add_text_box(slide, title, 0.45, y_pos + 0.04, 3.0, 0.35,
                 font_size=13, bold=True, color=C_DARK_BLUE)
    add_text_box(slide, code, 3.5, y_pos + 0.04, 9.4, 0.85,
                 font_size=11, color=C_NEAR_BLACK, font_name="Courier New")
    y_pos += 1.04


# ─────────────────────────────────────────────
# SLIDE 7 — What is XAI?
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "Explainable AI (XAI) — Why It Matters",
                   "Turning predictions into biological insights")
add_slide_number(slide, 7)

add_text_box(slide,
             "A model that says \"this is a CD4 T cell\" is useful.\n"
             "A model that also says \"…because SELL, CCR7, and IL7R are highly expressed\" is scientifically valuable.",
             0.5, 1.45, 12.3, 0.95, font_size=17, italic=True,
             color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

add_rect(slide, 0.3, 2.5, 12.73, 0.04, fill_color=C_LIGHT_BLUE)

cols = [
    (0.3, "Black-box\nproblem", C_ORANGE,
     "Deep neural nets give accurate predictions but the internal reasoning is opaque. "
     "Professors, reviewers, and clinicians need justification."),
    (4.6, "Biological\nvalidation", C_MID_BLUE,
     "If top-attributed genes match known cell-type markers, we gain confidence that the model learned real biology — not spurious correlations."),
    (8.9, "Discovery\npotential", C_GREEN,
     "High-scoring genes that are NOT known markers may represent novel biology or batch effects, guiding future wet-lab experiments."),
]

for cx, title, color, body in cols:
    add_rect(slide, cx, 2.65, 4.0, 4.55, fill_color=C_LIGHT_BLUE,
             line_color=color, line_width_pt=2)
    add_text_box(slide, title, cx + 0.15, 2.75, 3.7, 0.7,
                 font_size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, body, cx + 0.15, 3.5, 3.7, 3.55,
                 font_size=14, color=C_NEAR_BLACK, word_wrap=True)

add_text_box(slide,
             "scTRaCT uses two complementary gradient-based methods: Integrated Gradients & GradientSHAP",
             0.5, 7.1, 12.3, 0.35, font_size=13, italic=True,
             color=C_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 8 — Integrated Gradients Theory
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "Integrated Gradients (IG) — Theory",
                   "Sundararajan et al., ICML 2017  |  Captum: IntegratedGradients")
add_slide_number(slide, 8)

add_text_box(slide, "Core Idea", 0.5, 1.45, 12.0, 0.4,
             font_size=17, bold=True, color=C_DARK_BLUE)
add_text_box(slide,
             "IG asks: how much does each input feature contribute to the model's output, "
             "relative to a neutral baseline (e.g., zero expression)? "
             "It integrates the gradients of the output w.r.t. each input along a straight path "
             "from the baseline to the actual input.",
             0.5, 1.85, 12.0, 0.95, font_size=15, color=C_NEAR_BLACK, word_wrap=True)

# Formula box
add_rect(slide, 0.5, 2.85, 12.0, 0.85, fill_color=RGBColor(0xEE, 0xF5, 0xFF),
         line_color=C_MID_BLUE, line_width_pt=1.5)
add_text_box(slide,
             "IG_i(x) = (x_i − x'_i)  ×  ∫₀¹  [ ∂F(x' + α(x−x')) / ∂x_i ]  dα",
             0.7, 2.92, 11.6, 0.65, font_size=17, bold=True,
             color=C_DARK_BLUE, align=PP_ALIGN.CENTER, font_name="Courier New")

bullets = [
    ("x  = actual cell input (log-normalized gene expression vector)", []),
    ("x' = baseline (zero vector → represents a 'silent' cell with no expression)", []),
    ("F  = model output (logit for the predicted cell type)", []),
    ("α  ∈ [0,1] parameterizes the interpolation from baseline to input", []),
    ("The integral is approximated using n_steps = 50 uniform steps (Riemann sum)", []),
    ("Result: a score per gene — positive means the gene pushed the prediction toward this class", []),
    ("We take the mean of absolute values across sampled cells to get a cell-type–level ranking", []),
]

add_bullet_slide_content(slide, bullets, 0.5, 3.8, 12.0, 3.5,
                         font_size=14, color=C_NEAR_BLACK)

# Completeness axiom note
add_rect(slide, 0.5, 7.05, 12.0, 0.38, fill_color=RGBColor(0xFC, 0xF3, 0xCF),
         line_color=C_ORANGE, line_width_pt=1)
add_text_box(slide,
             "Axiom: IG satisfies Completeness — the sum of all attributions equals F(x) − F(x').",
             0.65, 7.08, 11.7, 0.32, font_size=12, color=C_ORANGE)


# ─────────────────────────────────────────────
# SLIDE 9 — GradientSHAP Theory
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "GradientSHAP — Theory",
                   "Lundberg & Lee (NIPS 2017)  |  Captum: GradientShap")
add_slide_number(slide, 9)

add_text_box(slide, "Core Idea", 0.5, 1.45, 12.0, 0.4,
             font_size=17, bold=True, color=C_DARK_BLUE)
add_text_box(slide,
             "GradientSHAP approximates SHAP (SHapley Additive exPlanations) values using gradient information. "
             "Instead of a single zero baseline, it uses a distribution of reference samples (100 randomly selected cells). "
             "Gradients are computed at noisy interpolations between the input and each baseline sample, then averaged.",
             0.5, 1.85, 12.0, 1.1, font_size=15, color=C_NEAR_BLACK, word_wrap=True)

# Comparison table
headers = ["Property", "Integrated Gradients", "GradientSHAP"]
rows = [
    ["Baseline",         "Single zero vector",                    "Distribution of real cells (100 random)"],
    ["Randomness",       "Deterministic",                         "Stochastic (noise + random baseline selection)"],
    ["Game theory link", "Path-based axioms",                     "Approximates Shapley values"],
    ["Computational cost","Medium (n_steps forward passes)",      "Lower (one pass per baseline sample)"],
    ["Best for",         "Precise per-gene attribution",          "Capturing natural data variation as reference"],
]

col_w = [3.0, 4.5, 4.5]
col_x = [0.5, 3.55, 8.1]
row_h = 0.55
start_y = 3.1

# Header row
for ci, (hdr, cw, cx) in enumerate(zip(headers, col_w, col_x)):
    add_rect(slide, cx, start_y, cw, row_h, fill_color=C_DARK_BLUE)
    add_text_box(slide, hdr, cx + 0.05, start_y + 0.08, cw - 0.1, row_h - 0.1,
                 font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    ry = start_y + row_h + ri * row_h
    fc = C_LIGHT_BLUE if ri % 2 == 0 else C_WHITE
    for ci, (cell_text, cw, cx) in enumerate(zip(row, col_w, col_x)):
        add_rect(slide, cx, ry, cw, row_h, fill_color=fc,
                 line_color=C_MID_BLUE, line_width_pt=0.5)
        add_text_box(slide, cell_text, cx + 0.05, ry + 0.05, cw - 0.1, row_h - 0.05,
                     font_size=12, color=C_NEAR_BLACK, word_wrap=True)

add_text_box(slide,
             "scTRaCT runs BOTH methods and returns IG_Score and SHAP_Score columns — "
             "genes ranking high in both are high-confidence markers.",
             0.5, 6.85, 12.0, 0.55, font_size=14, italic=True,
             color=C_MID_BLUE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 10 — XAI Pipeline Step-by-Step
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "scTRaCT XAI Pipeline — Step by Step",
                   "interpretability.py  |  explain_celltype()  /  explain_all_celltypes()")
add_slide_number(slide, 10)

steps_xai = [
    ("1", "Filter Cells",
     "Select all cells whose predicted label (obs['predicted_celltypes']) matches the target cell type.\n"
     "Sample up to num_cells=50 (or fewer if not available)."),
    ("2", "Prepare Tensors",
     "Convert lognorm layer → float32 tensor (counts).\n"
     "Compute inverse distance:  dist = 1 / (distance_matrix + 1e-6)  → float32 tensor."),
    ("3", "Run Integrated Gradients",
     "baseline = zero tensor (represents zero expression / silent cell).\n"
     "Integrate gradients along 50 steps from baseline → actual input.\n"
     "Extract count-stream attributions:  attr_ig[0]  (shape: cells × genes)."),
    ("4", "Run GradientSHAP",
     "background = 100 randomly sampled cells from the full adata.\n"
     "GradientSHAP computes gradient at noisy interpolations between each input and a random background sample.\n"
     "Extract count-stream attributions:  attr_shap[0]."),
    ("5", "Aggregate",
     "For both methods:  mean(|attributions|, axis=0)  → one score per gene.\n"
     "Build DataFrame: gene | IG_Score | SHAP_Score — sorted descending."),
    ("6", "Store & Visualize",
     "Saved to adata.uns['gene_attributions'][cell_type].\n"
     "Optional: horizontal bar charts (top N genes) + summary heatmap across all cell types."),
]

x_left = 0.3
y_start = 1.5
step_h = 0.88

for num, title, body in steps_xai:
    # number bubble
    add_rect(slide, x_left, y_start, 0.55, step_h - 0.05,
             fill_color=C_DARK_BLUE, line_color=C_DARK_BLUE)
    add_text_box(slide, num, x_left, y_start + 0.12, 0.55, 0.5,
                 font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # title
    add_text_box(slide, title, x_left + 0.65, y_start + 0.04, 2.4, 0.4,
                 font_size=13, bold=True, color=C_DARK_BLUE)
    # body
    add_text_box(slide, body, x_left + 3.15, y_start + 0.04, 9.8, step_h - 0.12,
                 font_size=12, color=C_NEAR_BLACK, word_wrap=True)
    # thin separator
    if num != "6":
        add_rect(slide, x_left + 0.65, y_start + step_h - 0.04,
                 12.35, 0.02, fill_color=C_LIGHT_BLUE)
    y_start += step_h


# ─────────────────────────────────────────────
# SLIDE 11 — Code Snippet
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "Using the XAI Module — Code Example")
add_slide_number(slide, 11)

code = """from scTRaCT import (prepare_data, TransformerModel, FocalLoss,
                     train_model, evaluate_on_query,
                     explain_celltype, explain_all_celltypes)

# ── 1. Prepare data ──────────────────────────────────────────────────────────
X_tr_c, X_tr_d, y_tr, X_val_c, X_val_d, y_val, le = prepare_data(adata)

# ── 2. Build & train model ───────────────────────────────────────────────────
model = TransformerModel(num_genes=X_tr_c.shape[1], num_classes=len(le.classes_))
train_model(model, criterion, X_tr_c, X_tr_d, y_tr, X_val_c, X_val_d, y_val)

# ── 3. Evaluate on query data ────────────────────────────────────────────────
acc, f1, preds, query_adata = evaluate_on_query(model_path, adata_query, le)
query_adata.obs['predicted_celltypes'] = preds   # required for XAI

# ── 4a. Explain a single cell type ───────────────────────────────────────────
df = explain_celltype(
    model, query_adata, le,
    cell_type="CD4 T",           # which cell type to explain
    method="both",               # IG + GradientSHAP
    num_cells=50,                # cells to sample
    n_steps=50,                  # IG integration steps
    plot=True                    # show bar chart
)

# ── 4b. Explain ALL predicted cell types at once ─────────────────────────────
results = explain_all_celltypes(
    model, query_adata, le,
    method="both",
    save_dir="attribution_plots/"    # saves PNGs + CSVs automatically
)
# Results also stored in: query_adata.uns['gene_attributions']"""

add_rect(slide, 0.3, 1.45, 12.73, 5.75, fill_color=RGBColor(0x1E, 0x1E, 0x2E),
         line_color=C_MID_BLUE, line_width_pt=1)
add_text_box(slide, code, 0.45, 1.5, 12.43, 5.65,
             font_size=10.5, color=RGBColor(0xCC, 0xFF, 0xCC),
             font_name="Courier New", word_wrap=False)


# ─────────────────────────────────────────────
# SLIDE 12 — Outputs & Visualizations
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "XAI Outputs & Visualizations")
add_slide_number(slide, 12)

# Left column — bar chart description
add_rect(slide, 0.3, 1.5, 5.9, 5.6, fill_color=C_LIGHT_BLUE,
         line_color=C_MID_BLUE, line_width_pt=1.5)
add_text_box(slide, "Per-Cell-Type Bar Chart",
             0.5, 1.6, 5.5, 0.45, font_size=15, bold=True, color=C_DARK_BLUE)
add_bullet_slide_content(slide, [
    "Horizontal bar chart — top N genes ranked by mean |attribution| score",
    "One chart per method (IG, SHAP)",
    "X-axis: Mean |Score| (always non-negative)",
    "Longer bar = gene had larger influence on the prediction",
    "Saved as:  <CellType>_attributions_IG.png",
    "Immediately interpretable by domain experts",
], 0.5, 2.1, 5.5, 4.7, font_size=13, color=C_NEAR_BLACK)

# Right column — heatmap description
add_rect(slide, 6.8, 1.5, 6.2, 5.6, fill_color=RGBColor(0xFD, 0xF2, 0xE9),
         line_color=C_ORANGE, line_width_pt=1.5)
add_text_box(slide, "Summary Heatmap (All Cell Types)",
             7.0, 1.6, 5.8, 0.45, font_size=15, bold=True, color=C_ORANGE)
add_bullet_slide_content(slide, [
    "Rows = predicted cell types",
    "Columns = union of top-N genes across all types",
    "Color = mean |attribution| score",
    "Reveals which genes are type-specific vs shared markers",
    "Saved as:  summary_heatmap_IG.png  /  summary_heatmap_SHAP.png",
    "Ideal for presentations and publications",
], 7.0, 2.1, 5.8, 4.7, font_size=13, color=C_NEAR_BLACK)

add_text_box(slide,
             "CSV files are also saved per cell type:  <CellType>_attributions.csv  "
             "— full gene ranking for downstream analysis",
             0.3, 7.05, 12.73, 0.38, font_size=12, italic=True,
             color=C_GRAY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 13 — Key Implementation Details
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "Key Implementation Details & Design Choices")
add_slide_number(slide, 13)

details = [
    ("Why mean of ABSOLUTE values?",
     "Genes can push a prediction either toward or away from a class. "
     "Taking the absolute value captures 'relevance' regardless of direction, "
     "then averaging across cells gives a stable cell-type–level ranking."),
    ("Why the COUNT stream only?",
     "The model has two input streams (counts + distances). Attribution is computed on both, "
     "but we report only the count-stream scores because they map 1-to-1 to interpretable gene names. "
     "Distance stream attributions are in the same gene space but reflect geometric proximity, not expression."),
    ("Why not use attention weights?",
     "The model has only 2 sequence tokens: [CLS] and [gene_embedding]. "
     "There is no per-gene attention weight to extract — the gene dimension is already compressed "
     "into a single vector before the Transformer. Gradient methods are therefore the correct tool here."),
    ("Why a random background for SHAP?",
     "SHAP expects a reference distribution that represents 'typical' data. "
     "Using 100 randomly sampled real cells from adata is more biologically realistic "
     "than a single zero vector, giving SHAP a more meaningful reference point."),
    ("Device handling",
     "Auto-detected at runtime (CUDA if available, else CPU). "
     "Never passed by the user — consistent with trainer.py."),
]

y = 1.5
for title, body in details:
    add_rect(slide, 0.3, y, 12.73, 1.0, fill_color=C_LIGHT_BLUE,
             line_color=C_MID_BLUE, line_width_pt=1)
    add_text_box(slide, title, 0.45, y + 0.05, 3.8, 0.4,
                 font_size=13, bold=True, color=C_DARK_BLUE)
    add_text_box(slide, body, 4.3, y + 0.05, 8.6, 0.88,
                 font_size=12, color=C_NEAR_BLACK, word_wrap=True)
    y += 1.06


# ─────────────────────────────────────────────
# SLIDE 14 — Comparison Table (IG vs SHAP vs Attention)
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "Method Comparison: IG vs GradientSHAP vs Attention")
add_slide_number(slide, 14)

headers2 = ["Criterion", "Integrated Gradients", "GradientSHAP", "Attention Weights"]
rows2 = [
    ["Theoretical grounding",  "Path-based axioms (completeness, sensitivity)", "Shapley values (fair contribution)", "None — attention ≠ explanation"],
    ["Baseline / reference",   "Single zero vector",                            "Distribution of real cells",          "N/A"],
    ["Deterministic?",         "Yes",                                            "No (random noise + baseline)",         "Yes"],
    ["Captures interactions?", "Partially (via gradient path)",                 "Better (via averaging over baselines)","No"],
    ["Used in scTRaCT?",       "Yes ✓",                                          "Yes ✓",                                "No — not applicable (2-token sequence)"],
    ["Output",                 "Gene scores per class",                          "Gene scores per class",                "Token attention weights"],
]

col_w2 = [3.2, 3.1, 3.1, 3.1]
col_x2 = [0.3, 3.55, 6.7, 9.85]
row_h2 = 0.6
start_y2 = 1.55

for ci, (hdr, cw, cx) in enumerate(zip(headers2, col_w2, col_x2)):
    add_rect(slide, cx, start_y2, cw, row_h2, fill_color=C_DARK_BLUE)
    add_text_box(slide, hdr, cx + 0.05, start_y2 + 0.1, cw - 0.1, row_h2 - 0.1,
                 font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for ri, row2 in enumerate(rows2):
    ry = start_y2 + row_h2 + ri * row_h2
    fc = C_LIGHT_BLUE if ri % 2 == 0 else C_WHITE
    for ci, (cell_t, cw, cx) in enumerate(zip(row2, col_w2, col_x2)):
        tc = C_GREEN if "✓" in cell_t else (C_ORANGE if "No" in cell_t and ci > 0 else C_NEAR_BLACK)
        add_rect(slide, cx, ry, cw, row_h2, fill_color=fc,
                 line_color=C_MID_BLUE, line_width_pt=0.5)
        add_text_box(slide, cell_t, cx + 0.05, ry + 0.05, cw - 0.1, row_h2 - 0.1,
                     font_size=11, color=tc, word_wrap=True)


# ─────────────────────────────────────────────
# SLIDE 15 — Summary & Future Work
# ─────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_WHITE)
add_section_header(slide, "Summary & Future Directions")
add_slide_number(slide, 15)

# Left: Summary
add_rect(slide, 0.3, 1.5, 6.0, 5.6, fill_color=C_LIGHT_BLUE,
         line_color=C_DARK_BLUE, line_width_pt=1.5)
add_text_box(slide, "What scTRaCT delivers",
             0.5, 1.6, 5.6, 0.45, font_size=15, bold=True, color=C_DARK_BLUE)
add_bullet_slide_content(slide, [
    "Accurate cell type classification using a Transformer trained on both gene expression and MCA distance features",
    "Handles class imbalance via Focal Loss — rare cell types are not ignored",
    "Two XAI methods (IG + GradientSHAP) identify which genes drove each prediction",
    "Results stored in AnnData — compatible with the Scanpy ecosystem",
    "Automated plots (bar charts + heatmap) and CSVs ready for publication",
], 0.5, 2.1, 5.7, 4.7, font_size=13, color=C_NEAR_BLACK)

# Right: Future work
add_rect(slide, 6.8, 1.5, 6.2, 5.6, fill_color=RGBColor(0xFD, 0xF2, 0xE9),
         line_color=C_ORANGE, line_width_pt=1.5)
add_text_box(slide, "Future Directions",
             7.0, 1.6, 5.8, 0.45, font_size=15, bold=True, color=C_ORANGE)
add_bullet_slide_content(slide, [
    "Cross-dataset benchmarking against Seurat, scANVI, and TOSICA",
    "Multi-dataset / cross-species transfer learning",
    "Integrate RNA velocity features as a third input stream",
    "Interactive visualization dashboard (Streamlit / Dash)",
    "Attribution-guided marker discovery pipeline for wet-lab validation",
    "Extend to spatial transcriptomics (spatial coordinates as additional context)",
], 7.0, 2.1, 5.8, 4.7, font_size=13, color=C_NEAR_BLACK)

add_text_box(slide,
             "Thank you!  Questions?",
             0.3, 7.08, 12.73, 0.38,
             font_size=18, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────

out_path = "/mnt/c/Users/mlw043/OneDrive - University of Texas at San Antonio/DOCTORAL PROGRAM AT UTSA/RESEARCH STUDY/scTRaCT/Python_Package/scTRaCT_XAI_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
